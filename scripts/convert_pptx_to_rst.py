#!/usr/bin/env python3
"""
Convert a .pptx file into a reStructuredText (.rst) document suitable for Sphinx/Read the Docs.

- Extracts slide titles, bullet points, and basic text content
- Exports embedded images and references them in the generated .rst

Usage:
  python scripts/convert_pptx_to_rst.py \
      --input "ML-life-science-Introduction to Jetstream2.pptx" \
      --output "docs/section1/ml_life_science_introduction_to_jetstream2.rst" \
      --media-dir "docs/section1/images/js2_pptx"

Requires:
  - python-pptx
  - Pillow (for robust image extension handling, optional but recommended)

Note:
  This is a best-effort textual conversion. Exact slide layouts and complex visuals
  may not be faithfully replicated. Embedded images are exported and referenced.
"""

import argparse
import os
import re
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    raise SystemExit(
        "Missing dependency: python-pptx. Install with `pip install python-pptx`"
    )

# Some environments may not have Pillow; we'll proceed without if unavailable
try:
    from PIL import Image  # noqa: F401
    PIL_AVAILABLE = True
except Exception:
    PIL_AVAILABLE = False


def slugify(text: str) -> str:
    text = text.strip().lower()
    text = re.sub(r"[^a-z0-9]+", "-", text)
    text = re.sub(r"-+", "-", text)
    return text.strip("-") or "slide"


def underline(text: str, char: str) -> str:
    return f"{text}\n{char * len(text)}\n\n"


def extract_text_from_shape(shape) -> list:
    lines = []
    if not hasattr(shape, "has_text_frame"):
        return lines
    if not shape.has_text_frame:
        return lines
    tf = shape.text_frame
    for p in tf.paragraphs:
        # Determine bullet level and text content
        level = p.level or 0
        text = "".join(run.text for run in p.runs).strip()
        if not text:
            continue
        indent = "  " * level
        # Use dash bullets; we don't attempt to infer ordered lists reliably here
        lines.append(f"{indent}- {text}")
    return lines


def image_extension_from_content_type(content_type: str) -> str:
    mapping = {
        "image/jpeg": "jpg",
        "image/png": "png",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tiff",
        "image/x-wmf": "wmf",
        "image/x-emf": "emf",
    }
    return mapping.get(content_type, "png")


def save_picture(shape, media_dir: Path, slide_idx: int, pic_idx: int) -> str:
    image = shape.image
    # Prefer extension from pptx image object if available
    ext = getattr(image, "ext", None)
    if not ext:
        ext = image_extension_from_content_type(getattr(image, "content_type", ""))
    filename = media_dir / f"slide{slide_idx+1}_img{pic_idx+1}.{ext}"
    with open(filename, "wb") as f:
        f.write(image.blob)
    # Return relative path from docs root for RST
    return str(filename)


def convert_pptx_to_rst(input_path: Path, output_rst: Path, media_dir: Path) -> None:
    prs = Presentation(str(input_path))

    media_dir.mkdir(parents=True, exist_ok=True)
    output_rst.parent.mkdir(parents=True, exist_ok=True)

    # Try to get a document title from first slide title
    doc_title = input_path.stem.replace("_", " ")
    if prs.slides:
        first = prs.slides[0]
        for shp in first.shapes:
            if shp.has_text_frame and getattr(shp, "is_placeholder", False):
                text = shp.text.strip()
                if text:
                    doc_title = text
                    break

    lines = []
    # File title
    lines.append(underline(doc_title, "="))
    lines.append("Note: This page was generated from a PowerPoint file. Formatting has been simplified.\n\n")

    # Iterate slides
    for s_idx, slide in enumerate(prs.slides):
        # Slide title
        slide_title = None
        if getattr(slide, "shapes", None):
            for shp in slide.shapes:
                if getattr(shp, "is_placeholder", False) and getattr(shp, "placeholder_format", None):
                    # Title placeholder has type 1 typically, but we'll accept first placeholder with text
                    if hasattr(shp, "has_text_frame") and shp.has_text_frame:
                        txt = shp.text.strip()
                        if txt:
                            slide_title = txt
                            break
        if not slide_title:
            slide_title = f"Slide {s_idx+1}"

        lines.append(underline(slide_title, "-"))

        # Text content and images
        pic_counter = 0
        collected_any_text = False
        for shp in slide.shapes:
            # Pictures
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                rel_path = save_picture(shp, media_dir, s_idx, pic_counter)
                pic_counter += 1
                # Use a relative path from the RST file location
                # Compute a relative path from output_rst.parent to rel_path
                rel_for_rst = os.path.relpath(rel_path, start=output_rst.parent)
                lines.append(f".. image:: {rel_for_rst}")
                lines.append("   :align: center")
                lines.append("   :width: 800px\n")
                continue

            # Text frames
            text_lines = extract_text_from_shape(shp)
            if text_lines:
                collected_any_text = True
                lines.extend(text_lines)

        if collected_any_text:
            lines.append("")

    output_rst.write_text("\n".join(lines), encoding="utf-8")
    print(f"Wrote RST: {output_rst}")


def main():
    parser = argparse.ArgumentParser(description="Convert PPTX to RST for Sphinx")
    parser.add_argument("--input", required=True, help="Path to input .pptx file")
    parser.add_argument(
        "--output",
        required=True,
        help="Path to output .rst file (e.g., docs/section1/my_page.rst)",
    )
    parser.add_argument(
        "--media-dir",
        required=True,
        help="Directory to write extracted images (e.g., docs/section1/images/js2_pptx)",
    )
    args = parser.parse_args()

    input_path = Path(args.input).expanduser().resolve()
    output_rst = Path(args.output).expanduser().resolve()
    media_dir = Path(args.media_dir).expanduser().resolve()

    if not input_path.exists():
        raise SystemExit(f"Input PPTX not found: {input_path}")

    convert_pptx_to_rst(input_path, output_rst, media_dir)


if __name__ == "__main__":
    main()
